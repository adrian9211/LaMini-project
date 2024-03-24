import base64
import streamlit as st
import torch
import streamlit_shadcn_ui as ui
from pptx import Presentation
from transformers import T5ForConditionalGeneration
from transformers import T5Tokenizer
from transformers import pipeline
from transformers import PegasusTokenizer, PegasusForConditionalGeneration
from sentence_transformers import SentenceTransformer, util



# MODEL and TOKENIZER
checkpoint = "LaMini-Flan-T5-248M"
tokenizer = T5Tokenizer.from_pretrained(checkpoint)
base_model = T5ForConditionalGeneration.from_pretrained(checkpoint, device_map='auto', torch_dtype=torch.float32)

#File Loader and Preprocessing
def file_processing(file):
  presentation = Presentation(file)
  text = []
  for slide in presentation.slides:
    for shape in slide.shapes:
      if hasattr(shape, "text"):
        text.append(shape.text)
  return "\n".join(text)

def extract_text_from_ppt(ppt_file):
  presentation = Presentation(ppt_file)
  text = []
  for slide in presentation.slides:
    for shape in slide.shapes:
      if hasattr(shape, "text"):
        text.append(shape.text)
  return "\n".join(text)

#LLM pipeline
def LLM_Pipeline(filepath):
  pipeline_summarizer = pipeline(
    'summarization',
    model = base_model,
    tokenizer = tokenizer,
    max_length = 152,
    min_length = 25
  )
  input_pptx = file_processing(filepath)
  result = pipeline_summarizer(input_pptx)
  result = result[0]['summary_text']
  return result


def summarize_and_display_result():
  uploaded_file = st.file_uploader("Upload your PPTX File", type=['pptx'])
  if uploaded_file is not None:
    if st.button("Summarize"):

      filepath = "data/"+uploaded_file.name
      with open(filepath, 'wb') as temp_file:
        temp_file.write(uploaded_file.read())


        st.info("Here is your PPTX Summarization")
        summary = LLM_Pipeline(filepath)
        st.success(summary)

def all_data_models_summarizer():
  # selected_data_model = data_models()
  #
  # model = SentenceTransformer(selected_data_model)
  def extract_text_from_ppt(ppt_file):
    presentation = Presentation(ppt_file)
    text = []
    for slide in presentation.slides:
      for shape in slide.shapes:
        if hasattr(shape, "text"):
          text.append(shape.text)
    return "\n".join(text)

  def summarise_text(text):
    summary = summariser(
      text,
      max_length=300,
      min_length=25,
      length_penalty=2.0,
      num_beams=4,
      early_stopping=True,
      do_sample=False
    )
    return summary[0]['summary_text']

  def find_relevant_text_by_topic(ppt_file, topics):
    presentation = Presentation(ppt_file)
    slide_texts = []
    for slide in presentation.slides:
      slide_text = " ".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
      slide_texts.append(slide_text)

    # Generate embeddings for all slide texts
    slide_embeddings = model.encode(slide_texts, convert_to_tensor=True)

    topic_texts = {}
    for topic in topics:
      # Generate embedding for the topic
      topic_embedding = model.encode([topic], convert_to_tensor=True)

      # Calculate similarities
      similarities = util.pytorch_cos_sim(topic_embedding, slide_embeddings)[0]

      # Find the slide with the highest similarity to the topic
      highest_similarity_index = torch.argmax(similarities).item()
      highest_similarity_score = similarities[highest_similarity_index].item()

      if highest_similarity_score > 0.3:  # ToDO: adjust based on the similarity requirement
        topic_texts[topic] = slide_texts[highest_similarity_index]
      else:
        topic_texts[topic] = "Relevant text for this topic was not found in the presentation."

    return topic_texts

  uploaded_file = st.file_uploader("Choose a PPT file", type="pptx")
  if uploaded_file is not None:
    # ToDO: improve the keywords lists or add a similarity model to extract similar content wrt to the topic
    topic_keywords = {
      "Problem Statement": ["problem", "issue", "challenge"],
      "Solution": ["solution", "resolve", "approach", "how we solve"],
      "Market Opportunity": ["market", "opportunity", "potential", "target audience"],
      "Business Model": ["business model", "revenue model", "how we make money"],
      "Unique Selling Proposition (USP)": ["unique selling proposition", "USP", "differentiators",
                                           "competitive advantage"],
      "Traction and Milestones": ["traction", "milestones", "progress", "achievements"],
      "Go-to-Market Strategy": ["go-to-market", "market strategy", "launch strategy"],
      "Team": ["team", "founders", "management", "leadership"],
      "Financial Projections": ["financial projections", "financials", "revenue forecast", "profitability"],
      "Investment Ask and Use of Funds": ["investment", "funding ask", "use of funds", "capital requirement"],
    }

    text = extract_text_from_ppt(uploaded_file)
    topic_texts = find_relevant_text_by_topic(uploaded_file, topic_keywords)

    tab1, tab2 = st.tabs(["Full Document Summary", "Topic Summary"])

    with tab1:

      # Summarise the full document
      # Model selection
      model_options = ["sshleifer/distilbart-cnn-12-6", "google/pegasus-xsum", "t5-small",
                       "facebook/bart-large-cnn"]
      model_choice = st.selectbox("Select Model", model_options)
      summariser = pipeline("summarization", model=model_choice)
      full_summary = summarise_text(text)
      st.write(full_summary)

    with tab2:

      # Topic based summarisation
      topic = st.selectbox('Choose a topic to summarise', ['Select a topic...'] + list(topic_keywords.keys()))

      if st.button('Generate Topic Summary'):
        if topic != 'Select a topic...':
          topic_text = topic_texts.get(topic, "This topic was not found in the presentation.")
          if "was not found" in topic_text:
            st.write(topic_text)
          else:
            topic_summary = summarise_text(topic_text)
            st.write(topic_summary)
def data_models():
    selected_data_model = st.selectbox("Choose Data Model", ["all-MiniLM-L6-v2",
                                                             "LaMini-Flan-T5-248M",
                                                             "all-distilroberta-v1",
                                                             "paraphrase-MiniLM-L6-v2",
                                                             "paraphrase-albert-small-v2",
                                                             "sentence-transformers/all-MiniLM-L12-v2",
                                                             "all-mpnet-base-v2"])
    return selected_data_model
def main():
  st.title('Legal-Pythia')
  st.subheader('Online PPTX Summarizer')

  selected_data_model= data_models()

  if selected_data_model == "LaMini-Flan-T5-248M":
    st.write("You have selected LaMini-Flan-T5-248M")
    summarize_and_display_result()
  elif selected_data_model == "all-MiniLM-L6-v2":
    st.write("You have selected all-MiniLM-L6-v2")

    # summarize_and_display_result()
    all_data_models_summarizer()
  elif selected_data_model == "all-distilroberta-v1":
    st.write("You have selected all-distilroberta-v1")
    all_data_models_summarizer()
  elif selected_data_model == "paraphrase-MiniLM-L6-v2":
    st.write("You have selected paraphrase-MiniLM-L6-v2")
    all_data_models_summarizer()
  elif selected_data_model == "paraphrase-albert-small-v2":
    st.write("You have selected paraphrase-albert-small-v2")
    all_data_models_summarizer()
  elif selected_data_model == "sentence-transformers/all-MiniLM-L12-v2":
    st.write("You have selected sentence-transformers/all-MiniLM-L12-v2")
    all_data_models_summarizer()
  elif selected_data_model == "all-mpnet-base-v2":
    st.write("You have selected all-mpnet-base-v2")
    all_data_models_summarizer()




if __name__ == "__main__":
  main()